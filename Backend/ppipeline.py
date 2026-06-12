

import io
import json
import mimetypes
import os
import re
import struct
import uuid
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import fitz  # PyMuPDF
import openpyxl
import structlog
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from app.storage.minio_client import MinIOClient

logger = structlog.get_logger(__name__)

minio = MinIOClient()


@dataclass
class ProcessedFile:
    file_id: str
    filename: str
    content_type: str
    extracted_text: str
    structured_content: Optional[Dict[str, Any]]
    metadata: Dict[str, Any]
    page_count: Optional[int]
    word_count: int
    processing_method: str  # "pymupdf" | "docx" | "openpyxl" | "pptx" | "plain_text"
    processing_status: str  # "complete" | "partial" | "failed"
    error: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return self.__dict__.copy()


class FileProcessingPipeline:
    """
    Main file processing pipeline.
    Routes files to the correct processor based on MIME type.
    """

    PROCESSORS = {
        "application/pdf": "process_pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "process_docx",
        "application/msword": "process_docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "process_xlsx",
        "application/vnd.ms-excel": "process_xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "process_pptx",
        "application/vnd.ms-powerpoint": "process_pptx",
        "text/plain": "process_text",
        "text/markdown": "process_text",
        "text/csv": "process_csv",
        "image/png": "process_image",
        "image/jpeg": "process_image",
        "image/jpg": "process_image",
        "image/webp": "process_image",
        "image/gif": "process_image",
    }

    async def process(
        self,
        storage_path: str,
        filename: str,
        content_type: str,
    ) -> ProcessedFile:
        """
        Main entry point. Downloads file from MinIO and processes it.
        """
        logger.info("Processing file", filename=filename, content_type=content_type)

        # Download from MinIO
        file_bytes = await minio.get_bytes(storage_path)
        if not file_bytes:
            return ProcessedFile(
                file_id=str(uuid.uuid4()),
                filename=filename,
                content_type=content_type,
                extracted_text="",
                structured_content=None,
                metadata={},
                page_count=None,
                word_count=0,
                processing_method="failed",
                processing_status="failed",
                error="File not found in storage",
            )

        # Route to correct processor
        processor_name = self.PROCESSORS.get(content_type)
        if not processor_name:
            # Try to infer from extension
            ext = Path(filename).suffix.lower()
            ext_map = {
                ".pdf": "process_pdf",
                ".docx": "process_docx",
                ".doc": "process_docx",
                ".xlsx": "process_xlsx",
                ".xls": "process_xlsx",
                ".pptx": "process_pptx",
                ".ppt": "process_pptx",
                ".txt": "process_text",
                ".md": "process_text",
                ".csv": "process_csv",
                ".png": "process_image",
                ".jpg": "process_image",
                ".jpeg": "process_image",
            }
            processor_name = ext_map.get(ext, "process_unknown")

        processor = getattr(self, processor_name, self.process_unknown)

        try:
            result = await processor(file_bytes, filename, content_type)
            return result
        except Exception as e:
            logger.error("File processing failed", filename=filename, error=str(e))
            return ProcessedFile(
                file_id=str(uuid.uuid4()),
                filename=filename,
                content_type=content_type,
                extracted_text="",
                structured_content=None,
                metadata={},
                page_count=None,
                word_count=0,
                processing_method="failed",
                processing_status="failed",
                error=str(e),
            )

    # ------------------------------------------------------------------
    # PDF processor
    # ------------------------------------------------------------------

    async def process_pdf(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        """
        Process PDF using PyMuPDF (fitz).
        Handles digital PDFs and attempts text extraction.
        Large PDFs are chunked; OCR available for scanned docs.
        """
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages_text: List[str] = []
        structured = {"pages": []}
        metadata = dict(doc.metadata) if doc.metadata else {}

        total_chars = 0
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            total_chars += len(text)

            # Extract links from page
            links = [{"uri": link.get("uri"), "rect": list(link["from"])} for link in page.get_links() if link.get("uri")]

            page_data = {
                "page_number": page_num + 1,
                "text": text,
                "char_count": len(text),
                "links": links,
            }
            pages_text.append(text)
            structured["pages"].append(page_data)

        # If minimal text extracted, mark for potential OCR
        processing_method = "pymupdf"
        processing_status = "complete"
        if total_chars < 100 and len(doc) > 0:
            logger.warning("PDF has minimal text — may be scanned", filename=filename)
            processing_status = "partial"
            processing_method = "pymupdf_low_text"

        full_text = "\n\n".join(filter(None, pages_text))

        # Extract table of contents
        toc = doc.get_toc()
        if toc:
            structured["toc"] = [{"level": t[0], "title": t[1], "page": t[2]} for t in toc]

        doc.close()

        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=full_text,
            structured_content=structured,
            metadata={
                **metadata,
                "page_count": len(pages_text),
                "total_chars": total_chars,
            },
            page_count=len(pages_text),
            word_count=len(full_text.split()),
            processing_method=processing_method,
            processing_status=processing_status,
        )

    # ------------------------------------------------------------------
    # DOCX processor
    # ------------------------------------------------------------------

    async def process_docx(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        """Process Word documents extracting text, headings, and structure."""
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
        except Exception as e:
            raise RuntimeError(f"Failed to open DOCX: {e}")

        paragraphs = []
        headings = []
        structured = {"sections": [], "tables": []}
        current_section = {"heading": None, "content": []}

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else "Normal"

            if style_name.startswith("Heading"):
                if current_section["content"]:
                    structured["sections"].append(current_section)
                level = int(style_name.split()[-1]) if style_name[-1].isdigit() else 1
                headings.append({"level": level, "text": text})
                current_section = {"heading": {"level": level, "text": text}, "content": []}
            else:
                paragraphs.append(text)
                current_section["content"].append(text)

        if current_section["content"]:
            structured["sections"].append(current_section)

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            table_data = {"index": table_idx, "rows": []}
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data["rows"].append(row_data)
            structured["tables"].append(table_data)

        # Core properties
        props = doc.core_properties
        metadata = {
            "author": props.author,
            "title": props.title,
            "subject": props.subject,
            "description": props.description,
            "keywords": props.keywords,
            "created": str(props.created) if props.created else None,
            "modified": str(props.modified) if props.modified else None,
            "revision": props.revision,
        }

        full_text = "\n".join(paragraphs)

        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=full_text,
            structured_content=structured,
            metadata=metadata,
            page_count=None,  # Word doesn't have fixed pages
            word_count=len(full_text.split()),
            processing_method="python-docx",
            processing_status="complete",
        )

    # ------------------------------------------------------------------
    # XLSX processor
    # ------------------------------------------------------------------

    async def process_xlsx(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        """Process Excel spreadsheets extracting all worksheets and tables."""
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)

        all_text: List[str] = []
        structured = {"worksheets": []}
        metadata = {
            "sheet_count": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data = {"name": sheet_name, "rows": [], "row_count": 0, "col_count": 0}
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                # Skip completely empty rows
                if not any(cell is not None for cell in row):
                    continue
                row_strs = [str(cell) if cell is not None else "" for cell in row]
                sheet_data["rows"].append(row_strs)
                all_text.append("\t".join(row_strs))
                row_count += 1

                # Cap rows per sheet to avoid memory issues
                if row_count > 10000:
                    all_text.append("[... truncated ...]")
                    break

            sheet_data["row_count"] = row_count
            structured["worksheets"].append(sheet_data)

        wb.close()
        full_text = "\n".join(all_text)

        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=full_text,
            structured_content=structured,
            metadata=metadata,
            page_count=len(wb.sheetnames),
            word_count=len(full_text.split()),
            processing_method="openpyxl",
            processing_status="complete",
        )

    # ------------------------------------------------------------------
    # PPTX processor
    # ------------------------------------------------------------------

    async def process_pptx(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        """Process PowerPoint presentations extracting slides, titles, notes, and text."""
        prs = Presentation(io.BytesIO(file_bytes))
        slides_data = []
        all_text: List[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_info = {
                "slide_number": slide_num,
                "title": None,
                "text_blocks": [],
                "notes": None,
                "shapes": [],
            }

            for shape in slide.shapes:
                # Extract title
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    slide_info["shapes"].append({"type": "image", "name": shape.name})
                    continue

                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue

                    # Check if this is a title placeholder
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        ph_type = shape.placeholder_format.type
                        if ph_type in (1, 13):  # CENTER_TITLE or TITLE
                            slide_info["title"] = text
                            all_text.append(f"## Slide {slide_num}: {text}")
                            continue

                    slide_info["text_blocks"].append(text)
                    all_text.append(text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_info["notes"] = notes_text
                        all_text.append(f"[Notes: {notes_text}]")

            slides_data.append(slide_info)

        structured = {
            "slide_count": len(prs.slides),
            "slides": slides_data,
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
        }

        full_text = "\n".join(all_text)

        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=full_text,
            structured_content=structured,
            metadata={"slide_count": len(prs.slides)},
            page_count=len(prs.slides),
            word_count=len(full_text.split()),
            processing_method="python-pptx",
            processing_status="complete",
        )

    # ------------------------------------------------------------------
    # Plain text / CSV processors
    # ------------------------------------------------------------------

    async def process_text(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        text = file_bytes.decode("utf-8", errors="replace")
        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=text,
            structured_content=None,
            metadata={"encoding": "utf-8"},
            page_count=1,
            word_count=len(text.split()),
            processing_method="plain_text",
            processing_status="complete",
        )

    async def process_csv(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        import csv
        text = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        structured = {
            "headers": rows[0] if rows else [],
            "rows": rows[1:] if len(rows) > 1 else [],
            "row_count": len(rows),
        }
        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text=text,
            structured_content=structured,
            metadata={"row_count": len(rows)},
            page_count=1,
            word_count=len(text.split()),
            processing_method="csv",
            processing_status="complete",
        )

    # ------------------------------------------------------------------
    # Image processor
    # ------------------------------------------------------------------

    async def process_image(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        """
        Process images: extract EXIF metadata.
        NOTE: We do NOT embed image content automatically.
        The user must confirm before we trigger vision-based extraction.
        """
        metadata = {}
        try:
            import struct
            # Basic PNG metadata
            if content_type == "image/png" and file_bytes[:8] == b'\x89PNG\r\n\x1a\n':
                metadata["format"] = "PNG"
                # Extract width/height from IHDR chunk
                metadata["width"] = struct.unpack(">I", file_bytes[16:20])[0]
                metadata["height"] = struct.unpack(">I", file_bytes[20:24])[0]
        except Exception:
            pass

        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text="",  # No text until user confirms OCR/vision processing
            structured_content=None,
            metadata={
                **metadata,
                "size_bytes": len(file_bytes),
                "awaiting_user_confirmation": True,
                "processing_note": "Image content extraction requires user confirmation before embedding.",
            },
            page_count=1,
            word_count=0,
            processing_method="image_metadata_only",
            processing_status="partial",
        )

    async def process_unknown(
        self, file_bytes: bytes, filename: str, content_type: str
    ) -> ProcessedFile:
        logger.warning("Unknown file type", filename=filename, content_type=content_type)
        return ProcessedFile(
            file_id=str(uuid.uuid4()),
            filename=filename,
            content_type=content_type,
            extracted_text="",
            structured_content=None,
            metadata={"size_bytes": len(file_bytes)},
            page_count=None,
            word_count=0,
            processing_method="unsupported",
            processing_status="failed",
            error=f"Unsupported file type: {content_type}",
        )
